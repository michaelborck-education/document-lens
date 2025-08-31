"""
Document processing service for text extraction from various file formats
"""

import json
import io
from typing import BinaryIO, Optional
from fastapi import HTTPException

# Import libraries for different file types
try:
    import PyPDF2
    from pdfplumber import PDF as PDFPlumber
except ImportError:
    PyPDF2 = None
    PDFPlumber = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

class DocumentProcessor:
    """Handles text extraction from various document formats"""
    
    async def extract_text(self, content: bytes, content_type: str, filename: str) -> str:
        """
        Extract text from document based on content type
        
        Args:
            content: Raw file content as bytes
            content_type: MIME type of the file
            filename: Name of the file (for error reporting)
            
        Returns:
            Extracted text as string
        """
        try:
            if content_type == "application/pdf":
                return await self._extract_from_pdf(content, filename)
            elif content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                return await self._extract_from_docx(content, filename)
            elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                return await self._extract_from_pptx(content, filename)
            elif content_type in ["text/plain", "text/markdown"]:
                return await self._extract_from_text(content, filename)
            elif content_type == "application/json":
                return await self._extract_from_json(content, filename)
            else:
                raise ValueError(f"Unsupported content type: {content_type}")
                
        except Exception as e:
            raise HTTPException(
                status_code=422,
                detail=f"Failed to extract text from {filename}: {str(e)}"
            )
    
    async def _extract_from_pdf(self, content: bytes, filename: str) -> str:
        """Extract text from PDF file"""
        if not PyPDF2:
            raise ImportError("PyPDF2 not available. Please install with: pip install PyPDF2")
        
        text_parts = []
        
        try:
            # Try with PyPDF2 first
            pdf_file = io.BytesIO(content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            for page in pdf_reader.pages:
                text = page.extract_text()
                if text.strip():
                    text_parts.append(text)
            
            # If PyPDF2 didn't extract much text, try pdfplumber
            if len("".join(text_parts)) < 100 and PDFPlumber:
                pdf_file.seek(0)
                with PDFPlumber(pdf_file) as pdf:
                    text_parts = []
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            text_parts.append(text)
            
        except Exception as e:
            raise ValueError(f"Failed to parse PDF: {str(e)}")
        
        if not text_parts:
            raise ValueError("No text could be extracted from PDF")
        
        return "\n\n".join(text_parts)
    
    async def _extract_from_docx(self, content: bytes, filename: str) -> str:
        """Extract text from DOCX file"""
        if not DocxDocument:
            raise ImportError("python-docx not available. Please install with: pip install python-docx")
        
        try:
            doc_file = io.BytesIO(content)
            doc = DocxDocument(doc_file)
            
            text_parts = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_parts.append(cell.text)
            
        except Exception as e:
            raise ValueError(f"Failed to parse DOCX: {str(e)}")
        
        if not text_parts:
            raise ValueError("No text could be extracted from DOCX")
        
        return "\n\n".join(text_parts)
    
    async def _extract_from_pptx(self, content: bytes, filename: str) -> str:
        """Extract text from PPTX file"""
        if not Presentation:
            raise ImportError("python-pptx not available. Please install with: pip install python-pptx")
        
        try:
            pptx_file = io.BytesIO(content)
            prs = Presentation(pptx_file)
            
            text_parts = []
            for slide in prs.slides:
                slide_texts = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)
                
                # Extract text from tables
                if hasattr(shape, "table"):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                slide_texts.append(cell.text)
                
                if slide_texts:
                    text_parts.append("\n".join(slide_texts))
            
        except Exception as e:
            raise ValueError(f"Failed to parse PPTX: {str(e)}")
        
        if not text_parts:
            raise ValueError("No text could be extracted from PPTX")
        
        return "\n\n".join(text_parts)
    
    async def _extract_from_text(self, content: bytes, filename: str) -> str:
        """Extract text from plain text or markdown file"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    text = content.decode(encoding)
                    return text
                except UnicodeDecodeError:
                    continue
            
            # If all encodings fail, use utf-8 with error handling
            return content.decode('utf-8', errors='replace')
            
        except Exception as e:
            raise ValueError(f"Failed to decode text file: {str(e)}")
    
    async def _extract_from_json(self, content: bytes, filename: str) -> str:
        """Extract text from JSON file"""
        try:
            text = content.decode('utf-8')
            data = json.loads(text)
            
            # Extract text recursively from JSON structure
            extracted_texts = []
            self._extract_text_from_json_recursive(data, extracted_texts)
            
            if not extracted_texts:
                raise ValueError("No text content found in JSON")
            
            return "\n\n".join(extracted_texts)
            
        except json.JSONDecodeError as e:
            raise ValueError(f"Invalid JSON format: {str(e)}")
        except Exception as e:
            raise ValueError(f"Failed to process JSON: {str(e)}")
    
    def _extract_text_from_json_recursive(self, obj, texts: list):
        """Recursively extract text from JSON object"""
        if isinstance(obj, dict):
            for key, value in obj.items():
                if isinstance(value, str) and len(value.strip()) > 10:
                    texts.append(value.strip())
                elif isinstance(value, (dict, list)):
                    self._extract_text_from_json_recursive(value, texts)
        elif isinstance(obj, list):
            for item in obj:
                if isinstance(item, str) and len(item.strip()) > 10:
                    texts.append(item.strip())
                elif isinstance(item, (dict, list)):
                    self._extract_text_from_json_recursive(item, texts)
    
    def validate_file_size(self, content: bytes, max_size: int) -> bool:
        """Validate file size"""
        return len(content) <= max_size